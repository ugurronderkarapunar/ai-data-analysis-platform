"""Automatic reporting agent with multi-format export support."""

from __future__ import annotations

from datetime import datetime
from pathlib import Path
from typing import Any, Dict, List, Optional

from core.base_agent import BaseAgent
from core.context import AgentContext
from config.settings import get_settings
from utils.exceptions import ReportError, ValidationError


class ReportingAgent(BaseAgent):
    """Generate executive / technical reports from pipeline outputs."""

    name = "reporting"
    description = "Yönetici ve teknik rapor üretir (Markdown/PDF/PPTX yapıda)"

    def run(self, context: AgentContext, **kwargs: Any) -> AgentContext:
        """Build report sections and optionally write Markdown to disk.

        Args:
            context: Shared context.
            **kwargs: Optional ``user_query``, ``write_files`` (bool).

        Returns:
            AgentContext: Context with report payload.

        Raises:
            ValidationError: If required analysis outputs are missing.
        """
        if context.dataframe is None:
            raise ValidationError("Rapor için dataframe gerekli.")

        user_query = str(kwargs.get("user_query") or "")
        write_files = bool(kwargs.get("write_files", True))

        try:
            report = self._compose(context, user_query)
            context.report = report
            if write_files:
                md_path = self._write_markdown(report)
                report["artifacts"] = {
                    "markdown": str(md_path),
                    "pdf_ready": True,
                    "pptx_ready": True,
                    "note": (
                        "PDF ve PowerPoint export fonksiyonları "
                        "report['markdown'] içeriğinden üretilebilir."
                    ),
                }
            context.add_message("Otomatik rapor oluşturuldu.")
            self.logger.info("Report generated")
            return context
        except ValidationError:
            raise
        except Exception as exc:  # noqa: BLE001
            self.logger.exception("Reporting failed")
            raise ReportError(f"Rapor oluşturulamadı: {exc}") from exc

    def _compose(self, context: AgentContext, user_query: str) -> Dict[str, Any]:
        """Compose structured report sections.

        Args:
            context: Shared context.
            user_query: Original user query.

        Returns:
            Dict[str, Any]: Report dictionary.
        """
        profile = context.profile or {}
        stats = context.statistics or {}
        forecast = context.forecast or {}
        ml = context.ml_advice or {}
        cleaning = context.cleaning_plan or {}

        executive = [
            f"Kaynak dosya: {context.file_name or 'bilinmiyor'}",
            f"Kapsam: {profile.get('rows', '?')} satır, {profile.get('columns', '?')} sütun.",
            f"Kullanıcı isteği: {user_query or 'Genel analiz'}",
            f"Birincil niyet: {(context.intent or {}).get('primary_intent', 'general_analysis')}",
        ]
        if forecast.get("available"):
            executive.append(
                f"Tahmin modeli: {forecast.get('selected_model')} "
                f"(sıklık: {forecast.get('frequency')})."
            )
        if ml.get("recommendations"):
            executive.append(
                f"Önerilen ML yaklaşımı: {ml['recommendations'][0]['type']}."
            )

        technical = [
            f"Eksik hücre sayısı: {cleaning.get('missing_total', profile.get('missing_counts') and sum(profile.get('missing_counts', {}).values()) or 0)}",
            f"Yinelenen satır: {profile.get('duplicate_rows', 0)}",
            f"Bellek: {profile.get('memory_mb', 0)} MB",
            f"Tarih sütunları: {', '.join(profile.get('datetime_columns', [])) or 'yok'}",
            f"Olası hedef alanlar: {', '.join(profile.get('target_candidates', [])[:5]) or 'yok'}",
        ]
        for note in stats.get("plain_language", [])[:6]:
            technical.append(note)

        insights: List[str] = []
        strong = (profile.get("correlations") or {}).get("strong_pairs") or []
        for pair in strong[:3]:
            insights.append(
                f"Güçlü korelasyon: {pair['a']} ↔ {pair['b']} (r={pair['corr']})."
            )
        for rec in (ml.get("recommendations") or [])[:3]:
            insights.append(f"{rec['type']}: {rec['why']}")
        if forecast.get("explanation"):
            insights.append(forecast["explanation"])
        if not insights:
            insights.append("Veri genel keşif için uygun; iş sorusu netleştirilerek derinleştirilebilir.")

        recommendations = [
            "Temizleme planındaki adımları iş kuralına göre onaylayın; otomatik uygulatmayın.",
            "Dashboard filtrelerini iş birimi KPI'larına göre daraltın.",
        ]
        if cleaning.get("recommendations"):
            recommendations.append(
                f"{len(cleaning['recommendations'])} temizleme önerisi incelenebilir."
            )
        if forecast.get("available"):
            recommendations.append(
                "Tahmin ufkunu iş takviminize göre (1H/1A/3A/6A/1Y) takip edin."
            )
        else:
            recommendations.append(
                "Zaman serisi tahmini için düzenli tarih + metrik kolonları sağlayın."
            )

        next_steps = [
            "İş sorusunu tek cümlelik KPI tanımına indirgeyin.",
            "Gerekirse SQL/Power BI/RAG ajanlarını registry üzerinden ekleyin.",
            "Onay sonrası temizlenmiş veri setini outputs/ altına sabitleyin.",
            "Model eğitimi gerekiyorsa ML danışmanı önerisini ürün sahibi ile doğrulayın.",
        ]

        markdown = self._to_markdown(
            executive,
            technical,
            insights,
            recommendations,
            next_steps,
            context,
        )
        return {
            "title": "Yapay Zekâ Veri Analiz Raporu",
            "generated_at": datetime.utcnow().isoformat() + "Z",
            "executive_summary": executive,
            "technical_summary": technical,
            "business_insights": insights,
            "recommendations": recommendations,
            "next_steps": next_steps,
            "markdown": markdown,
            "tasks": context.tasks,
        }

    def _to_markdown(
        self,
        executive: List[str],
        technical: List[str],
        insights: List[str],
        recommendations: List[str],
        next_steps: List[str],
        context: AgentContext,
    ) -> str:
        """Render report sections as Markdown.

        Args:
            executive: Executive bullets.
            technical: Technical bullets.
            insights: Insight bullets.
            recommendations: Recommendation bullets.
            next_steps: Next step bullets.
            context: Shared context.

        Returns:
            str: Markdown document.
        """
        lines = [
            "# Yapay Zekâ Veri Analiz Raporu",
            "",
            f"**Dosya:** {context.file_name or '-'}",
            f"**Üretim:** {datetime.utcnow().strftime('%Y-%m-%d %H:%M UTC')}",
            "",
            "## Yönetici Özeti",
            *[f"- {x}" for x in executive],
            "",
            "## Teknik Özet",
            *[f"- {x}" for x in technical],
            "",
            "## İş İçgörüleri",
            *[f"- {x}" for x in insights],
            "",
            "## Öneriler",
            *[f"- {x}" for x in recommendations],
            "",
            "## Sonraki Adımlar",
            *[f"- {x}" for x in next_steps],
            "",
        ]
        return "\n".join(lines)

    def _write_markdown(self, report: Dict[str, Any]) -> Path:
        """Persist Markdown report under outputs/.

        Args:
            report: Report payload.

        Returns:
            Path: Written file path.

        Raises:
            ReportError: If writing fails.
        """
        settings = get_settings()
        stamp = datetime.utcnow().strftime("%Y%m%d_%H%M%S")
        path = settings.outputs_dir / f"report_{stamp}.md"
        try:
            path.write_text(report["markdown"], encoding="utf-8")
            self.logger.info("Wrote markdown report to %s", path)
            return path
        except OSError as exc:
            raise ReportError(f"Markdown yazılamadı: {exc}") from exc


def export_pdf_from_markdown(markdown_text: str, output_path: Path) -> Path:
    """Export Markdown text to a simple PDF when reportlab is available.

    Args:
        markdown_text: Markdown content.
        output_path: Destination PDF path.

    Returns:
        Path: Output path.

    Raises:
        ReportError: If PDF backend is unavailable or write fails.
    """
    logger_local = ReportingAgent().logger
    try:
        from reportlab.lib.pagesizes import A4
        from reportlab.pdfgen import canvas
    except ImportError as exc:
        raise ReportError("PDF için reportlab kurulu değil.") from exc

    try:
        output_path.parent.mkdir(parents=True, exist_ok=True)
        c = canvas.Canvas(str(output_path), pagesize=A4)
        width, height = A4
        y = height - 50
        for raw in markdown_text.splitlines():
            line = raw[:110]
            if y < 50:
                c.showPage()
                y = height - 50
            c.drawString(40, y, line)
            y -= 14
        c.save()
        logger_local.info("PDF exported to %s", output_path)
        return output_path
    except Exception as exc:  # noqa: BLE001
        raise ReportError(f"PDF export başarısız: {exc}") from exc


def export_pptx_from_report(report: Dict[str, Any], output_path: Path) -> Path:
    """Export report sections to PowerPoint when python-pptx is available.

    Args:
        report: Report dictionary.
        output_path: Destination PPTX path.

    Returns:
        Path: Output path.

    Raises:
        ReportError: If PPTX backend is unavailable or write fails.
    """
    logger_local = ReportingAgent().logger
    try:
        from pptx import Presentation
    except ImportError as exc:
        raise ReportError("PowerPoint için python-pptx kurulu değil.") from exc

    try:
        prs = Presentation()
        sections = [
            ("Yönetici Özeti", report.get("executive_summary", [])),
            ("Teknik Özet", report.get("technical_summary", [])),
            ("İş İçgörüleri", report.get("business_insights", [])),
            ("Öneriler", report.get("recommendations", [])),
            ("Sonraki Adımlar", report.get("next_steps", [])),
        ]
        for title, bullets in sections:
            slide = prs.slides.add_slide(prs.slide_layouts[1])
            slide.shapes.title.text = title
            body = slide.shapes.placeholders[1].text_frame
            body.clear()
            for i, bullet in enumerate(bullets):
                if i == 0:
                    body.text = str(bullet)
                else:
                    body.add_paragraph().text = str(bullet)
        output_path.parent.mkdir(parents=True, exist_ok=True)
        prs.save(str(output_path))
        logger_local.info("PPTX exported to %s", output_path)
        return output_path
    except Exception as exc:  # noqa: BLE001
        raise ReportError(f"PPTX export başarısız: {exc}") from exc
